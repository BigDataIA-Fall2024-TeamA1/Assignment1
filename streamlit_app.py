# Import necessary libraries
import streamlit as st
from sql_module import (
    get_metadata_from_sql,
    update_metadata_steps,
    insert_evaluation,
    get_evaluations
)
from collections import Counter
from openai_module import send_to_openai
from aws_module import get_files_from_s3
import tempfile
import PyPDF2
import docx
import pandas as pd
import zipfile
import os
import openpyxl
import plotly.express as px
import boto3
import easyocr
import base64
import re

# Define a consistent color palette
COLOR_PALETTE = {
    'green': '#2ca02c',
    'red': '#d62728',
    'blue': '#1f77b4',
    'orange': '#ff7f0e',
    'purple': '#9467bd',
    'cyan': '#17becf',
    'grey': '#7f7f7f'
}

# Supported file types
SUPPORTED_EXTENSIONS = [
    '.json', '.pdf', '.png', '.jpeg', '.jpg', '.txt',
    '.xlsx', '.csv', '.pptx', '.docx', '.py', '.zip', '.pdb'
]

# Initialize EasyOCR reader
reader = easyocr.Reader(['en'], gpu=False)

# Function to extract text from different file types
def extract_text_from_file(file_path):
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()

    try:
        if ext == '.txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()

        elif ext == '.pdf':
            text = ""
            with open(file_path, 'rb') as f:
                reader_pdf = PyPDF2.PdfReader(f)
                for page in reader_pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n"
            return text

        elif ext == '.docx':
            doc = docx.Document(file_path)
            return "\n".join([para.text for para in doc.paragraphs])

        elif ext == '.csv':
            df = pd.read_csv(file_path)
            return df.to_csv(index=False)

        elif ext == '.xlsx':
            try:
                df = pd.read_excel(file_path, engine='openpyxl')
                csv_representation = df.to_csv(index=False)
                return csv_representation
            except Exception as e:
                st.error(f"Error processing Excel file: {e}")
                return f"Error processing Excel file: {e}"

        elif ext in ['.png', '.jpg', '.jpeg']:
            try:
                result = reader.readtext(file_path, detail=0)
                text = ' '.join(result)
                return text
            except Exception as e:
                st.error(f"Error processing image file: {e}")
                return f"Error processing image file: {e}"

        elif ext == '.py':
            return extract_text_from_py(file_path)

        elif ext == '.zip':
            return extract_text_from_zip(file_path)

        elif ext == '.pdb':
            return extract_text_from_pdb(file_path)

        elif ext == '.pptx':
            from pptx import Presentation
            prs = Presentation(file_path)
            text_runs = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_runs.append(shape.text)
            return "\n".join(text_runs)

        else:
            return f"Unsupported file type: {ext}"

    except Exception as e:
        return f"Error processing file: {e}"

# Helper function to extract text from .py files
def extract_text_from_py(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except UnicodeDecodeError:
        with open(file_path, 'r', encoding='latin-1') as f:
            return f.read()
    except Exception as e:
        return f"Error processing .py file: {e}"

# Helper function to extract text from .zip files
def extract_text_from_zip(file_path):
    text = ""
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        with tempfile.TemporaryDirectory() as tmpdir:
            zip_ref.extractall(tmpdir)
            for root, _, files in os.walk(tmpdir):
                for file in files:
                    file_path_inner = os.path.join(root, file)
                    _, ext_inner = os.path.splitext(file_path_inner)
                    ext_inner = ext_inner.lower()
                    if ext_inner in SUPPORTED_EXTENSIONS:
                        extracted_text = extract_text_from_file(file_path_inner)
                        text += f"Extracted from {file}:\n{extracted_text}\n"
                    else:
                        text += f"Skipped unsupported file: {file}\n"
    return text

# Helper function to extract text from .pdb files
def extract_text_from_pdb(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except Exception as e:
        return f"Error processing .pdb file: {e}"

# Initialize Streamlit session state
def initialize_session_state():
    session_keys = [
        'comparison_result',
        'openai_response',
        'final_answer',
        'steps',
        'selected_task_id',
        'show_feedback_form',
        'show_rerun_button',
        'modified_steps',
        'awaiting_rerun_satisfaction',
        'awaiting_feedback'
    ]
    for key in session_keys:
        if key not in st.session_state:
            if key in ['comparison_result', 'show_feedback_form', 'show_rerun_button', 'awaiting_rerun_satisfaction', 'awaiting_feedback']:
                st.session_state[key] = False
            else:
                st.session_state[key] = ""

initialize_session_state()

# Custom CSS for styling
st.markdown("""
    <style>
    /* Global styles */
    body {
        font-family: 'Arial', sans-serif;
        font-size: 18px;
        color: #ffffff;
        background-color: #000000;
    }

    /* Header styles */
    h1 {
        font-size: 40px;
        font-weight: bold;
        color: #ffffff;
        text-align: center;
    }

    h2, h3 {
        color: #ffffff;
    }

    /* Paragraph and div styles */
    p, div {
        font-size: 18px;
        color: #ffffff;
    }

    /* Table styling */
    table {
        width: 100%;
        border-collapse: collapse;
    }

    th, td {
        padding: 14px;
        text-align: left;
        border-bottom: 1px solid #444444;
        font-size: 16px;
        color: #ffffff;
    }

    th {
        background-color: #333333;
        color: #ffffff;
    }

    /* Button styling */
    .stButton > button {
        background-color: #444444;
        color: #ffffff;
        border-radius: 10px;
        border: 1px solid #ffffff;
    }

    /* Selectbox styling */
    .stSelectbox > div {
        background-color: #333333;
        color: #ffffff;
        border: 1px solid #ffffff;
        border-radius: 5px;
    }

    /* Dataframe styling */
    .stDataFrame {
        background-color: #333333;
        color: #ffffff;
        border: 1px solid #ffffff;
        border-radius: 5px;
    }

    /* Sidebar styling */
    .sidebar .sidebar-content {
        background-color: #333333;
    }

    /* App container background */
    .stApp {
        background-color: #000000;
    }
    </style>
    """, unsafe_allow_html=True)

# Streamlit App Layout

# Sidebar for navigation and branding
with st.sidebar:
    st.image("https://streamlit.io/images/brand/streamlit-logo-primary-colormark-darktext.png", width=150)
    st.header("Navigation")
    st.write("Use the main page to process Task IDs and evaluate OpenAI responses.")
    st.markdown("---")
    st.header("Settings")

# Main Page Header
st.markdown("""
    <h1>
        OpenAI Service Evaluation and Feedback Dashboard
    </h1>
    """, unsafe_allow_html=True)

st.markdown("---")

# Main Content Container
with st.container():
    # 1. Get metadata task_ids and questions from SQL Server
    metadata = get_metadata_from_sql()
    metadata_task_ids = [record['task_id'] for record in metadata]
    questions_dict = {record['task_id']: record['Question'] for record in metadata}
    final_answers_dict = {record['task_id']: record['Final answer'] for record in metadata}
    steps_dict = {record['task_id']: record['Steps'] for record in metadata}

    # 2. Get files from AWS S3 and create a mapping of task_ids to their files
    bucket_name = st.secrets["aws"]["bucket"]
    s3_files = get_files_from_s3(bucket_name)

    # Create a mapping from task_id to list of files (in case there are multiple files per task_id)
    task_files_mapping = {}
    for file_name in s3_files:
        file_base_name, file_ext = os.path.splitext
