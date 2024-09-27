import os
import boto3
from io import BytesIO
import pandas as pd
import docx
import zipfile
from PIL import Image
import pytesseract
import csv
import json
import wave
import jsonlines
import PyPDF2
from pptx import Presentation
import filetype
import tempfile

# Initialize S3 client
s3_client = boto3.client(
    's3',
    aws_access_key_id=os.getenv('AWS_ACCESS_KEY_ID'),
    aws_secret_access_key=os.getenv('AWS_SECRET_ACCESS_KEY')
)

# Download file from S3
def download_file_from_s3(bucket_name, file_key):
    """Download file from S3"""
    s3 = boto3.client('s3')
    file_obj = s3.get_object(Bucket=bucket_name, Key=file_key)
    return file_obj['Body'].read()

# Function to process PDF
def process_pdf(file_content):
    """Process PDF files"""
    reader = PyPDF2.PdfFileReader(BytesIO(file_content))
    text = ""
    for page_num in range(reader.numPages):
        page = reader.getPage(page_num)
        text += page.extractText()
    return text

# Function to process PPTX
def process_pptx(file_content):
    """Process PPTX files"""
    prs = Presentation(BytesIO(file_content))
    text = ""
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text += shape.text + "\n"
    return text

# Function to process TXT
def process_txt(file_content):
    """Process TXT files"""
    return file_content.decode('utf-8')

# Function to process DOCX
def process_docx(file_content):
    """Process DOCX files"""
    doc = docx.Document(BytesIO(file_content))
    text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
    return text

# Function to process CSV
def process_csv(file_content):
    """Process CSV files"""
    df = pd.read_csv(BytesIO(file_content))
    return df.to_string()

# Function to process JSON
def process_json(file_content):
    """Process JSON files"""
    data = json.loads(file_content)
    return json.dumps(data, indent=2)

# Function to process JSONL
def process_jsonl(file_content):
    """Process JSONL files"""
    with jsonlines.Reader(BytesIO(file_content)) as reader:
        return "\n".join([json.dumps(obj) for obj in reader])

# Function to process MP3 or WAV audio files (transcription placeholder)
def process_audio(file_content, file_extension):
    """Process MP3 or WAV audio files"""
    # This can be replaced with an actual transcription library like Whisper
    return f"Processed audio file with extension {file_extension}."

# Function to process ZIP
def process_zip(file_content):
    """Process ZIP files"""
    with tempfile.TemporaryDirectory() as temp_dir:
        with zipfile.ZipFile(BytesIO(file_content), 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
            extracted_files = os.listdir(temp_dir)
        return f"Extracted files: {', '.join(extracted_files)}"

# Function to process JPG/PNG using OCR (text recognition)
def process_image(file_content):
    """Process image files (JPG/PNG) using OCR"""
    with tempfile.NamedTemporaryFile(delete=False) as temp_img:
        temp_img.write(file_content)
        temp_img.flush()
        text = pytesseract.image_to_string(Image.open(temp_img.name))
    return text

# Function to process PY files (display code)
def process_py(file_content):
    """Process Python (.py) files"""
    return file_content.decode('utf-8')

# Function to process PDB (placeholder for molecular data)
def process_pdb(file_content):
    """Process PDB files"""
    return file_content.decode('utf-8')

# File processor based on file extension
def process_file(file_content, file_extension):
    """Dispatch to the appropriate file processing function based on the extension"""
    if file_extension == '.pdf':
        return process_pdf(file_content)
    elif file_extension == '.pptx':
        return process_pptx(file_content)
    elif file_extension == '.txt':
        return process_txt(file_content)
    elif file_extension == '.docx':
        return process_docx(file_content)
    elif file_extension == '.csv':
        return process_csv(file_content)
    elif file_extension == '.json':
        return process_json(file_content)
    elif file_extension == '.jsonl':
        return process_jsonl(file_content)
    elif file_extension in ['.mp3', '.wav']:
        return process_audio(file_content, file_extension)
    elif file_extension == '.zip':
        return process_zip(file_content)
    elif file_extension in ['.jpg', '.jpeg', '.png']:
        return process_image(file_content)
    elif file_extension == '.py':
        return process_py(file_content)
    elif file_extension == '.pdb':
        return process_pdb(file_content)
    else:
        return "Unsupported file format."

# Function to handle file processing from S3
def process_file_from_s3(bucket_name, file_key):
    """Download file from S3 and process it based on its extension"""
    file_content = download_file_from_s3(bucket_name, file_key)
    file_extension = os.path.splitext(file_key)[1].lower()  # Get file extension
    return process_file(file_content, file_extension)